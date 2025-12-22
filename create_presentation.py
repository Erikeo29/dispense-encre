from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_methodology_presentation():
    """Crée une présentation PowerPoint avec 2 slides sur la méthodologie de modélisation"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE 1 ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Titre
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "MÉTHODOLOGIE DE MODÉLISATION - PHYSIQUE & ÉQUATIONS"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    title_p.alignment = PP_ALIGN.CENTER

    # Approche de simulation
    y_pos = 1.0
    approach_box = slide1.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
    approach_frame = approach_box.text_frame
    approach_frame.word_wrap = True

    p = approach_frame.paragraphs[0]
    p.text = "Approche de simulation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for text in [
        "• Simulation diphasique (encre Ag/AgCl + air) en domaine microfluidique",
        "• Méthode Phase-Field pour le suivi d'interface",
        "• Régime laminaire incompressible"
    ]:
        p = approach_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        p.space_before = Pt(3)

    # Équations gouvernantes
    y_pos = 2.0
    eq_box = slide1.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(3.8))
    eq_frame = eq_box.text_frame
    eq_frame.word_wrap = True

    p = eq_frame.paragraphs[0]
    p.text = "Équations gouvernantes"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    equations = [
        ("1. Navier-Stokes (conservation de la masse et quantité de mouvement)",
         "   ∇·v = 0",
         "   ρ(∂v/∂t + v·∇v) = -∇p + ∇·τ + ρg + F_σ"),

        ("2. Rhéologie : Modèle de Carreau (fluide rhéofluidifiant)",
         "   η(γ̇) = η_∞ + (η₀ - η_∞)[1 + (λγ̇)²]^((n-1)/2)"),

        ("3. Transport d'interface (Phase-Field)",
         "   ∂φ/∂t + v·∇φ = γ∇·[ε∇φ - φ(1-φ²)n]"),

        ("4. Tension de surface",
         "   F_σ = σκδ(φ)n")
    ]

    for eq_group in equations:
        p = eq_frame.add_paragraph()
        p.text = eq_group[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.space_before = Pt(8)

        for eq_line in eq_group[1:]:
            p = eq_frame.add_paragraph()
            p.text = eq_line
            p.font.size = Pt(10)
            p.font.name = "Courier New"
            p.space_before = Pt(2)

    # Conditions aux limites
    y_pos = 6.0
    bc_box = slide1.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(1.0))
    bc_frame = bc_box.text_frame
    bc_frame.word_wrap = True

    p = bc_frame.paragraphs[0]
    p.text = "Conditions aux limites clés"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    for text in [
        "• Parois : Non-glissement (v = 0)",
        "• Angle de contact : n_w·∇φ = -(1/ε)cos(θ)",
        "• Entrée : v = v₀ (0.1 m/s), φ = 1 (encre)",
        "• Sortie : p = p_atm"
    ]:
        p = bc_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.space_before = Pt(3)

    # ==================== SLIDE 2 ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "DONNÉES D'ENTRÉE & VARIABLES DE SORTIE"
    title_p2 = title_frame2.paragraphs[0]
    title_p2.font.size = Pt(24)
    title_p2.font.bold = True
    title_p2.font.color.rgb = RGBColor(0, 51, 102)
    title_p2.alignment = PP_ALIGN.CENTER

    # Données d'entrée
    y_pos = 1.0
    input_box = slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(4.2))
    input_frame = input_box.text_frame
    input_frame.word_wrap = True

    p = input_frame.paragraphs[0]
    p.text = "Données d'entrée"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    input_data = [
        ("Géométrie", [
            "Diamètre puit (D_w): 800 - 1500 μm",
            "Diamètre seringue (D_s): 200 - 350 μm",
            "Décalage horizontal (Δx): 0, -75, -150 μm",
            "Hauteur puit (h_w): 128 μm (fixe)"
        ]),
        ("Rhéologie", [
            "Viscosité au repos (η₀): 0.5 - 5 Pa·s",
            "Viscosité cisaillement infini (η_∞): 0.05 Pa·s",
            "Temps relaxation (λ): 0.15 s",
            "Indice pseudoplasticité (n): 0.7"
        ]),
        ("Mouillage", [
            "Angle contact électrode or (θ_or): 35 - 70°",
            "Angle contact paroi (θ_wall): 35 - 90°"
        ]),
        ("Interface", [
            "Tension de surface (σ): 40 mN/m"
        ]),
        ("Process", [
            "Temps dispense: 40 ms",
            "Vitesse injection (v₀): 0.1 m/s"
        ])
    ]

    for category, items in input_data:
        p = input_frame.add_paragraph()
        p.text = f"• {category}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.space_before = Pt(8)

        for item in items:
            p = input_frame.add_paragraph()
            p.text = f"   - {item}"
            p.font.size = Pt(10)
            p.space_before = Pt(2)

    # Variables de sortie
    y_pos = 5.5
    output_box = slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(1.5))
    output_frame = output_box.text_frame
    output_frame.word_wrap = True

    p = output_frame.paragraphs[0]
    p.text = "Variables de sortie"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    output_data = [
        ("Champs spatiaux-temporels", "Vitesse (u, v), Pression (p), Fraction volumique (φ)"),
        ("Dynamique interface", "Position interface, Courbure (κ), Taux cisaillement (γ̇)"),
        ("Performance dispense", "Taux remplissage (%), Symétrie, Mouillage parois"),
        ("Visualisation", "Animation GIF 2D de l'écoulement")
    ]

    for category, description in output_data:
        p = output_frame.add_paragraph()
        p.text = f"• {category}: {description}"
        p.font.size = Pt(11)
        p.space_before = Pt(5)

    # Sauvegarder la présentation
    output_file = "Methodologie_Modelisation_Dispense.pptx"
    prs.save(output_file)
    return output_file

if __name__ == "__main__":
    filename = create_methodology_presentation()
    print(f"✓ Présentation créée : {filename}")
